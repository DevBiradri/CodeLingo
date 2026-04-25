import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# --- NEOBRUTALIST DESIGN TOKENS ---
BLACK = RGBColor(0, 0, 0)
WHITE = RGBColor(255, 255, 255)
YELLOW = RGBColor(255, 230, 0)  # #FFE600
PINK = RGBColor(255, 0, 255)    # #FF00FF
CYAN = RGBColor(0, 255, 255)    # #00FFFF
GRAY = RGBColor(240, 240, 240)

SHADOW_OFFSET = Pt(6)
BORDER_WIDTH = Pt(4)

def add_neobrutalist_box(slide, left, top, width, height, color=WHITE, shadow=True):
    """Adds a box with a thick black border and a hard shadow."""
    if shadow:
        # Draw shadow first (offset rectangle)
        shadow_shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            left + SHADOW_OFFSET, top + SHADOW_OFFSET, width, height
        )
        shadow_shape.fill.solid()
        shadow_shape.fill.fore_color.rgb = BLACK
        shadow_shape.line.fill.background() # No line for shadow

    # Draw main box
    box = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        left, top, width, height
    )
    box.fill.solid()
    box.fill.fore_color.rgb = color
    box.line.color.rgb = BLACK
    box.line.width = BORDER_WIDTH
    return box

def set_text_style(text_frame, text, font_size=Pt(24), bold=True, color=BLACK, align=PP_ALIGN.LEFT):
    """Helper to style text in a frame with better compatibility."""
    text_frame.clear() # Clear any default paragraph
    p = text_frame.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = font_size
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = "Arial" # More compatible font

def create_pitch_deck():
    prs = Presentation()
    # Set slide dimensions to 16:9
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slides_data = [
        {
            "title": "CodeLingo",
            "subtitle": "Learn Coding by Doing.\nInteractive. Structured. Addictive.",
            "color": YELLOW,
            "type": "title"
        },
        {
            "title": "The Problem",
            "content": [
                "Learning coding feels overwhelming",
                "Passive tutorials don’t build real skill",
                "Beginners get no instant feedback",
                "Most users quit before building anything"
            ],
            "color": PINK,
            "image": "landing_screen.png"
        },
        {
            "title": "The Solution",
            "content": [
                "Gamified learning experience",
                "Hands-on interactive challenges",
                "Instant feedback on every step",
                "Start coding with zero setup"
            ],
            "color": CYAN,
            "image": "map_screen.png"
        },
        {
            "title": "The Core Loop",
            "content": [
                "Health: Mistakes have impact",
                "Streaks: Build daily consistency",
                "XP: Track your progress",
                "Leagues: Stay motivated with competition"
            ],
            "color": YELLOW,
            "image": "profile_screen.png"
        },
        {
            "title": "The Visual Roadmap",
            "content": [
                "Clear skill progression path",
                "Structured learning journey",
                "Unlock new concepts step by step",
                "Designed to reach real-world readiness"
            ],
            "color": WHITE,
            "image": "map_screen.png"
        },
        {
            "title": "Multi-Modal Learning",
            "content": [
                "Predict Output: Test understanding",
                "Drag & Drop: Build correct syntax",
                "Refactor: Improve existing code",
                "Editor: Practice real coding"
            ],
            "color": PINK,
            "image": "challenge_screen.png"
        },
        {
            "title": "AI Engine",
            "content": [
                "Dynamically generated challenges",
                "Smart evaluation of code",
                "Context-aware hints",
                "Personalized learning pace"
            ],
            "color": CYAN,
            "image": "challenge_screen.png"
        },
        {
            "title": "Technical Stack",
            "content": [
                "Frontend: Next.js + Tailwind",
                "Backend: FastAPI",
                "AI: Gemini API",
                "Editor: Monaco Editor"
            ],
            "color": WHITE,
            "image": "landing_screen.png"
        },
        {
            "title": "Future Vision",
            "content": [
                "Multiplayer coding challenges",
                "Real-world project modules",
                "Enterprise learning platform",
                "Skill-based progression system"
            ],
            "color": YELLOW,
            "image": "leaderboard_screen.png"
        },
        {
            "title": "Join CodeLingo",
            "subtitle": "Make coding simple.\nMake learning consistent.\nMake progress real.",
            "color": PINK,
            "type": "title"
        }
    ]

    # Use absolute path for assets directory relative to the script
    script_dir = os.path.dirname(os.path.abspath(__file__))
    assets_dir = os.path.join(script_dir, "assets")
    
    if not os.path.exists(assets_dir):
        os.makedirs(assets_dir)

    for i, data in enumerate(slides_data):
        print(f"Generating Slide {i+1}: {data['title']}...")
        # Create a blank slide
        slide_layout = prs.slide_layouts[6] # Blank layout
        slide = prs.slides.add_slide(slide_layout)
        
        # Set background color
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = GRAY

        if data.get("type") == "title":
            # --- TITLE SLIDE LAYOUT ---
            box_w, box_h = Inches(10), Inches(4)
            left, top = (prs.slide_width - box_w)/2, (prs.slide_height - box_h)/2
            
            box = add_neobrutalist_box(slide, left, top, box_w, box_h, data["color"])
            
            # Title Text
            set_text_style(box.text_frame, data["title"], font_size=Pt(80), align=PP_ALIGN.CENTER)
            
            # Subtitle Box
            sub_w, sub_h = Inches(8), Inches(1.5)
            sub_left, sub_top = (prs.slide_width - sub_w)/2, top + box_h + Inches(0.5)
            sub_box = add_neobrutalist_box(slide, sub_left, sub_top, sub_w, sub_h, WHITE)
            set_text_style(sub_box.text_frame, data["subtitle"], font_size=Pt(28), align=PP_ALIGN.CENTER)

        else:
            # --- CONTENT SLIDE LAYOUT ---
            # Title Header
            title_w, title_h = Inches(6), Inches(1)
            title_box = add_neobrutalist_box(slide, Inches(0.5), Inches(0.5), title_w, title_h, data["color"])
            set_text_style(title_box.text_frame, data["title"], font_size=Pt(40))

            # Content Box (Left)
            content_w, content_h = Inches(5.5), Inches(4.5)
            content_box = add_neobrutalist_box(slide, Inches(0.5), Inches(2), content_w, content_h, WHITE)
            tf = content_box.text_frame
            tf.clear() # Clear default
            
            for point in data["content"]:
                p = tf.add_paragraph()
                p.space_before = Pt(12)
                run = p.add_run()
                run.text = f"- {point}"
                run.font.size = Pt(24)
                run.font.bold = True
                run.font.name = "Arial"

            # Image Box (Right)
            img_w, img_h = Inches(6), Inches(4.5)
            img_left, img_top = Inches(6.5), Inches(2)
            
            img_path = os.path.join(assets_dir, data.get("image", ""))
            
            if os.path.exists(img_path):
                # Add a black border/box behind the image
                add_neobrutalist_box(slide, img_left, img_top, img_w, img_h, WHITE)
                slide.shapes.add_picture(img_path, img_left + Pt(2), img_top + Pt(2), width=img_w - Pt(4), height=img_h - Pt(4))
            else:
                # Placeholder box
                box = add_neobrutalist_box(slide, img_left, img_top, img_w, img_h, WHITE)
                set_text_style(box.text_frame, f"[ {data.get('image', 'No Image')} ]\nPlace screenshot in scripts/assets/", align=PP_ALIGN.CENTER)

    output_path = os.path.join(os.path.dirname(script_dir), "CodeLingo_Pitch_Deck.pptx")
    prs.save(output_path)
    print(f"\nSUCCESS! Neobrutalist Pitch Deck created at: {os.path.abspath(output_path)}")
    print(f"Tip: Remember to put your screenshots in {assets_dir} and run this again!")

if __name__ == "__main__":
    create_pitch_deck()
